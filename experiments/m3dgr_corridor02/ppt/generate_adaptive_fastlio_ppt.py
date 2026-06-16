import os
import warnings
from pathlib import Path

os.environ.setdefault(
    "MPLCONFIGDIR",
    "/home/romi/ros2_ws/experiments/m3dgr_corridor02/ppt/.matplotlib",
)

warnings.filterwarnings(
    "ignore",
    message="Unable to import Axes3D.*",
    category=UserWarning,
)

import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/romi/ros2_ws")
RESULT_DIR = ROOT / "experiments/m3dgr_corridor02/results"
OUT_DIR = ROOT / "experiments/m3dgr_corridor02/ppt"
OUT_FILE = OUT_DIR / "adaptive_fast_lio2_idea_and_experiment.pptx"
FORMULA_DIR = OUT_DIR / "formula_images"

TITLE = RGBColor(28, 54, 92)
ACCENT = RGBColor(36, 139, 194)
DARK = RGBColor(38, 38, 38)
MUTED = RGBColor(95, 95, 95)
PANEL_LINE = RGBColor(205, 218, 232)


def set_run(run, size=22, bold=False, color=DARK, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, 29, True, TITLE)

    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.1), Inches(12), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        set_run(r, 13, False, MUTED)


def add_bullets(slide, items, x=0.8, y=1.35, w=11.8, h=5.3, size=21):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(9)
        for r in p.runs:
            set_run(r, size, False, DARK)


def add_code_box(slide, text, x, y, w, h, size=16):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = PANEL_LINE

    box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    for r in p.runs:
        set_run(r, size, False, RGBColor(45, 45, 45), "Consolas")


def add_caption(slide, text, x, y, w):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, 13, True, TITLE)


def add_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def render_formula(name, formula, fontsize=28):
    FORMULA_DIR.mkdir(parents=True, exist_ok=True)
    out_path = FORMULA_DIR / f"{name}.png"

    fig = plt.figure(figsize=(8.0, 1.0), dpi=240)
    fig.patch.set_alpha(0.0)
    plt.axis("off")
    plt.text(
        0.5,
        0.5,
        formula,
        ha="center",
        va="center",
        fontsize=fontsize,
        color="#222222",
        math_fontfamily="cm",
    )
    plt.savefig(out_path, transparent=True, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    return out_path


def add_formula_card(slide, title, formulas, notes, x, y, w, h):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = PANEL_LINE

    title_box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.08), Inches(w - 0.28), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, 15, True, TITLE)

    formula_y = y + 0.58
    for name, formula, fw, fh in formulas:
        img = render_formula(name, formula)
        slide.shapes.add_picture(
            str(img),
            Inches(x + (w - fw) / 2.0),
            Inches(formula_y),
            width=Inches(fw),
            height=Inches(fh),
        )
        formula_y += fh + 0.18


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. title / module position
    slide = prs.slides.add_slide(blank)
    add_title(slide, "退化感知的 FAST-LIO2 地图增量更新模块", "重点：adaptive_map 模块设计与 M3DGR Corridor02 实验结果")
    add_bullets(slide, [
        "研究对象：FAST-LIO2 中当前帧点云插入 ikd-Tree 地图之前的 map_incremental() 阶段。",
        "核心问题：原始入图逻辑主要考虑体素冗余，缺少对残差异常、几何退化和重复方向点的判断。",
        "本文模块：利用 h_share_model() 已有的点到面残差 $r_i$、有效匹配点数量 $N_{eff}$ 和平面法向量 $\\mathbf{n}_i$，对候选入图点进行质量筛选。",
        "目标：在走廊等退化环境中减少低质量地图点累积，提高后续 scan-to-map 匹配可靠性。"
    ], x=0.85, y=1.45, w=11.7, h=4.6, size=21)

    # 2. method pipeline
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块实现流程", "从观测模型中提取质量信息，在地图更新前执行筛选")
    add_code_box(slide, "1. h_share_model()\n\n输入: 当前帧降采样点 + ikd-Tree 局部地图\n\n提取质量信息:\n- 点到面残差 $|r_i|$\n- 平面法向量 $\\mathbf{n}_i$\n- 质量分数 $s_i$\n- 有效匹配点数量 $N_{eff}$", 0.75, 1.45, 3.75, 4.45, 15)
    add_code_box(slide, "2. adaptive_map 判断\n\n帧级指标:\n- 平均残差 $\\bar r$\n- 残差中位数 $\\tilde r$\n- MAD 鲁棒统计\n- 法向量退化指标 $\\rho_n$\n\n作用: 判断当前帧是否存在退化风险。", 4.8, 1.45, 3.75, 4.45, 15)
    add_code_box(slide, "3. map_incremental()\n\n在插入 ikd-Tree 前执行:\n- 拒绝残差异常点\n- 拒绝低质量分数点\n- 限制退化方向重复点\n- 少量放行新区域点\n\n关闭 adaptive_map 时回到原始逻辑。", 8.85, 1.45, 3.75, 4.45, 15)

    # 3. scoring and degeneracy
    slide = prs.slides.add_slide(blank)
    add_title(slide, "核心筛选逻辑", "单点质量评价 + 帧级退化判断 + 法向量方向约束")
    add_formula_card(
        slide,
        "单点质量评价",
        [
            ("residual", r"$r_i=\mathbf{n}_i^{T}\mathbf{p}_i^w+d_i$", 3.45, 0.5),
            ("score", r"$s_i=1-0.9\frac{|r_i|}{\sqrt{\|\mathbf{p}_i^w\|}}$", 3.55, 0.62),
        ],
        ["r_i 越小，当前点与局部地图平面越一致。"],
        0.75,
        1.32,
        5.75,
        2.25,
    )
    add_formula_card(
        slide,
        "自适应残差阈值",
        [
            ("residual_limit", r"$r_{lim}=\min\left(\tau_p,\max\left(\alpha\bar r,\tilde r+k\cdot1.4826\cdot MAD\right)\right)$", 5.0, 0.78),
            ("mad", r"$MAD=\mathrm{median}\left(|r_i-\tilde r|\right)$", 2.9, 0.45),
        ],
        ["结合平均残差与鲁棒统计量，避免异常点支配阈值。"],
        6.85,
        1.32,
        5.75,
        2.25,
    )
    add_formula_card(
        slide,
        "帧级退化判断",
        [
            ("normal_matrix", r"$\mathbf{N}=\frac{1}{N_{eff}}\sum_i\mathbf{n}_i\mathbf{n}_i^T$", 3.35, 0.55),
            ("normal_ratio", r"$\rho_n=\frac{\lambda_{min}(\mathbf{N})}{\lambda_{max}(\mathbf{N})}$", 2.95, 0.55),
        ],
        ["法向量方向越集中，rho_n 越小。"],
        0.75,
        4.05,
        5.75,
        2.25,
    )
    add_formula_card(
        slide,
        "最终入图拒绝条件",
        [
            ("reject", r"$|r_i|>r_{lim}\quad \mathrm{or}\quad s_i<s_{min}$", 3.7, 0.5),
            ("degenerate", r"$N_{eff}<N_{min}\quad \mathrm{or}\quad \bar r>r_{max}\quad \mathrm{or}\quad \rho_n<\rho_{min}$", 5.05, 0.42),
        ],
        ["退化帧中进一步限制重复法向量方向的入图点数量。"],
        6.85,
        4.05,
        5.75,
        2.25,
    )

    # 4. experiment setup
    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验设置", "M3DGR Corridor02 走廊退化场景")
    add_code_box(slide, "数据集与场景\n\nM3DGR Corridor02\nLivox Mid360 + IMU\n\n退化特点:\n- 长走廊\n- 重复结构\n- 主几何约束方向单一\n- 适合验证退化感知地图更新", 0.9, 1.45, 5.4, 4.25, 17)
    add_code_box(slide, "对比组与观察内容\n\n1. FAST-LIO2 官方版本\n2. Adaptive FAST-LIO2, adaptive_map 开启\n3. Adaptive FAST-LIO2, adaptive_map 关闭\n\n观察重点:\n地图主体结构、边界清晰度、局部散点噪声。", 7.0, 1.45, 5.4, 4.25, 17)

    # 5. result comparison
    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验结果对比", "三种方法均恢复走廊主体结构，adaptive_map 开启后局部散点有所收敛")
    add_picture(slide, RESULT_DIR / "fastlio.png", 0.42, 1.5, 3.95, 4.45)
    add_picture(slide, RESULT_DIR / "adaptive_on.png", 4.69, 1.5, 3.95, 4.45)
    add_picture(slide, RESULT_DIR / "adaptive_off.png", 8.96, 1.5, 3.95, 4.45)
    add_caption(slide, "FAST-LIO2", 0.42, 6.08, 3.95)
    add_caption(slide, "Adaptive ON", 4.69, 6.08, 3.95)
    add_caption(slide, "Adaptive OFF", 8.96, 6.08, 3.95)

    # 6. conclusion
    slide = prs.slides.add_slide(blank)
    add_title(slide, "结果分析与结论", "该模块在不破坏主体结构的前提下，对低质量点累积具有抑制作用")
    add_bullets(slide, [
        "Adaptive OFF 与官方 FAST-LIO2 地图形态接近，说明复现系统在关闭 adaptive_map 后基本保持原始入图行为。",
        "Adaptive ON 保持了 Corridor02 的矩形走廊主体结构，说明质量筛选没有造成明显过度删点或地图断裂。",
        "与 Adaptive OFF 相比，Adaptive ON 的局部散点和低质量累积有所减少，符合退化感知地图更新的设计目标。",
        "下一步建议加入定量指标：入图点数量、拒绝点比例、平均残差、有效匹配点比例、运行时间，以及有真值时的 ATE/RPE。"
    ], size=21)

    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build()
